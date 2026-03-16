"""
Topic Database – Optimized Streamlit Application
=================================================
Features:
  - Full CRUD with auto-generated TopicID, DaysOpen, aging bucket, risk score, escalation flag
  - Dashboard: KPI cards, aging bar chart, category Pareto, high-risk agenda list
  - Analytics page: monthly trend, root-cause Pareto, cross-category comparison
  - Smart creation form: multi-file upload, links, auto-dates, auto-naming attachments
  - Storage: Excel (data/topics.xlsx) + attachments (data/attachments/) + backups
  - Exports: Excel, CSV, PPTX
"""

import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime, date, timedelta
import io, json, os, shutil, re
from pathlib import Path

# ── Optional deps (soft imports) ───────────────────────────────────────────────
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Config ──────────────────────────────────────────────────────────────────────
DATA_DIR        = Path("data")
EXCEL_FILE      = DATA_DIR / "topics.xlsx"
ATTACH_DIR      = DATA_DIR / "attachments"
BACKUP_DIR      = DATA_DIR / "backups"
LINKS_COL       = "links"          # stored as JSON list of {"label":…,"url":…}

CATEGORIES  = ["Elements", "Assembly", "Cross"]
SEVERITIES  = ["Low", "Medium", "High", "Critical"]
STATUSES    = ["Open", "In Progress", "Blocked", "Closed"]
BUCKET_ORDER= ["0–3 Months", "3–6 Months", "6–12 Months", "> 1 Year"]

CAT_HEX  = {"Elements": "#3B82F6", "Assembly": "#10B981", "Cross": "#F59E0B"}
STAT_HEX = {"Open": "#EF4444", "In Progress": "#3B82F6", "Blocked": "#F59E0B", "Closed": "#10B981"}
SEV_HEX  = {"Low": "#6EE7B7", "Medium": "#FCD34D", "High": "#FCA5A5", "Critical": "#F87171"}

for d in [DATA_DIR, ATTACH_DIR, BACKUP_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ── Page config ─────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Topic Database",
    page_icon="🗂️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS ──────────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500&family=Mulish:wght@300;400;500;600;700&display=swap');

html, body, [class*="css"] { font-family: 'Mulish', sans-serif; }
header[data-testid="stHeader"] { background: transparent; }

/* KPI cards */
.kpi-card {
    background: #FFFFFF; border: 1px solid #E5E7EB; border-radius: 14px;
    padding: 22px 20px; text-align: center;
    box-shadow: 0 2px 8px rgba(0,0,0,0.05); transition: all .2s;
    position: relative; overflow: hidden;
}
.kpi-card::before {
    content: ''; position: absolute; top: 0; left: 0; right: 0; height: 3px;
    background: var(--accent, #3B82F6);
}
.kpi-card:hover { box-shadow: 0 8px 24px rgba(0,0,0,0.10); transform: translateY(-2px); }
.kpi-label { font-family: 'Syne', sans-serif; font-size: 11px; font-weight: 700;
             color: #6B7280; text-transform: uppercase; letter-spacing: .1em; margin-bottom: 8px; }
.kpi-value { font-family: 'Syne', sans-serif; font-size: 36px; font-weight: 800;
             color: #111827; line-height: 1; }
.kpi-sub   { font-size: 11px; color: #9CA3AF; margin-top: 6px; font-family: 'JetBrains Mono', monospace; }

/* Section title */
.section-title {
    font-family: 'Syne', sans-serif; font-size: 22px; font-weight: 800; color: #111827;
    margin: 0 0 18px 0; padding-bottom: 10px; border-bottom: 3px solid #111827; letter-spacing: -.02em;
}

/* Progress */
.progress-outer { background:#F3F4F6; border-radius:99px; height:8px; overflow:hidden; margin-top:8px; }
.progress-inner { height:100%; border-radius:99px; transition: width .6s ease; }

/* Bucket bars */
.bucket-row   { display:flex; align-items:center; gap:10px; margin-bottom:10px; font-size:13px; }
.bucket-label { width:130px; color:#374151; flex-shrink:0; font-weight:600; font-size:12px; }
.bucket-bar-outer { flex:1; background:#F3F4F6; border-radius:99px; height:10px; overflow:hidden; }
.bucket-bar-inner { height:100%; border-radius:99px; transition: width .6s ease; }
.bucket-count { width:30px; text-align:right; color:#6B7280; font-size:12px;
                font-family:'JetBrains Mono',monospace; font-weight:500; }

/* Tag pills */
.tag-pill {
    display:inline-block; padding:2px 10px; border-radius:99px;
    font-size:11px; font-weight:700; margin:2px; text-transform:uppercase; letter-spacing:.05em;
}

/* Topic card */
.topic-card {
    border: 1px solid #E5E7EB; border-radius: 12px; padding: 16px 20px;
    margin-bottom: 12px; background: #FAFAFA;
    border-left: 4px solid var(--card-accent, #3B82F6);
    transition: box-shadow .2s;
}
.topic-card:hover { box-shadow: 0 4px 16px rgba(0,0,0,0.08); }

/* Risk badge */
.risk-badge {
    display:inline-block; padding:3px 10px; border-radius:6px;
    font-size:11px; font-weight:700; font-family:'JetBrains Mono',monospace;
}

/* Sidebar */
section[data-testid="stSidebar"] { background: #0F172A !important; }
section[data-testid="stSidebar"] * { color: #E2E8F0 !important; }
section[data-testid="stSidebar"] .stRadio label { color: #CBD5E1 !important; font-size:14px; }
section[data-testid="stSidebar"] hr { border-color: #334155 !important; }
</style>
""", unsafe_allow_html=True)


# ── Helper functions ─────────────────────────────────────────────────────────────

def aging_bucket(opening_date):
    if pd.isna(opening_date):
        return "Unknown"
    days = (datetime.today() - pd.to_datetime(opening_date)).days
    if days <= 90:    return "0–3 Months"
    elif days <= 180: return "3–6 Months"
    elif days <= 365: return "6–12 Months"
    else:             return "> 1 Year"

def days_open(opening_date, status, close_date=None):
    if pd.isna(opening_date):
        return 0
    if status == "Closed" and close_date and not pd.isna(close_date):
        return max(0, (pd.to_datetime(close_date) - pd.to_datetime(opening_date)).days)
    return max(0, (datetime.today() - pd.to_datetime(opening_date)).days)

def risk_score(row):
    """0–100 risk score based on severity, days open, customer impact, status."""
    sev_w   = {"Low": 10, "Medium": 30, "High": 60, "Critical": 100}
    stat_w  = {"Open": 1.0, "In Progress": 0.7, "Blocked": 1.2, "Closed": 0.0}
    d_open  = days_open(row.get("opening_date"), row.get("status","Open"), row.get("close_date"))
    base    = sev_w.get(row.get("severity", "Medium"), 30)
    ci_mult = 1.4 if row.get("customer_impact") else 1.0
    age_mult= 1.0 + min(d_open / 365.0, 1.5)
    s_mult  = stat_w.get(row.get("status", "Open"), 1.0)
    score   = base * ci_mult * age_mult * s_mult
    return min(100, round(score))

def escalation_flag(row):
    """True if risk > 70 or days open > 90 with status not Closed."""
    if row.get("status") == "Closed":
        return False
    rs   = risk_score(row)
    d    = days_open(row.get("opening_date"), row.get("status","Open"), row.get("close_date"))
    return rs > 70 or d > 90

def enrich(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df["_days_open"]   = df.apply(lambda r: days_open(r.get("opening_date"), r.get("status","Open"), r.get("close_date")), axis=1)
    df["_bucket"]      = df["opening_date"].apply(aging_bucket)
    df["_risk"]        = df.apply(risk_score, axis=1)
    df["_escalate"]    = df.apply(escalation_flag, axis=1)
    return df

def sanitize_filename(name: str) -> str:
    return re.sub(r'[^\w\-.]', '_', name)


# ── Data I/O ─────────────────────────────────────────────────────────────────────

COLUMNS = [
    "id", "topic_group", "sub_topic", "category", "severity",
    "opening_date", "close_date", "pic",
    "problem_description", "root_cause_analysis", "corrective_actions",
    "next_steps", "status", "customer_impact",
    "milestone_dates",   # JSON: [{"label":…,"date":…}]
    "links",             # JSON: [{"label":…,"url":…}]
    "attachments",       # JSON: ["filename", …]
]

DEMO_DATA = [
    {"id": 1, "topic_group": "Welding Defects", "sub_topic": "Porosity in MIG Welds",
     "category": "Elements", "severity": "High",
     "opening_date": datetime.today() - timedelta(days=20), "close_date": None,
     "pic": "J. Müller", "problem_description": "Porosity in batch WD-22; ~15% rejection rate.",
     "root_cause_analysis": "Shielding gas contamination detected on line 3.",
     "corrective_actions": "Replaced gas supply line; tightened all fittings.",
     "next_steps": "Requalification test scheduled for next week.",
     "status": "In Progress", "customer_impact": True,
     "milestone_dates": json.dumps([{"label": "8D Due", "date": str(date.today() + timedelta(days=7))}]),
     "links": json.dumps([{"label": "8D Report", "url": "https://example.com/8d/1"}]),
     "attachments": json.dumps([])},
    {"id": 2, "topic_group": "Assembly Sequence", "sub_topic": "Bolt torque out-of-spec",
     "category": "Assembly", "severity": "Medium",
     "opening_date": datetime.today() - timedelta(days=130), "close_date": datetime.today() - timedelta(days=10),
     "pic": "A. Schmidt", "problem_description": "Torque values 10% below spec on rear bracket.",
     "root_cause_analysis": "Calibration drift on station 7 torque wrench.",
     "corrective_actions": "Wrench recalibrated; process audit completed.",
     "next_steps": "", "status": "Closed", "customer_impact": False,
     "milestone_dates": json.dumps([]), "links": json.dumps([]), "attachments": json.dumps([])},
    {"id": 3, "topic_group": "Supplier Quality", "sub_topic": "Dimensional deviation – Part X401",
     "category": "Cross", "severity": "Critical",
     "opening_date": datetime.today() - timedelta(days=250), "close_date": None,
     "pic": "L. Bauer", "problem_description": "OD of X401 exceeds tolerance +0.3 mm consistently.",
     "root_cause_analysis": "Under investigation – supplier audit planned.",
     "corrective_actions": "Interim: 100% incoming inspection.",
     "next_steps": "Awaiting supplier response to 8D.", "status": "Blocked", "customer_impact": True,
     "milestone_dates": json.dumps([{"label": "Supplier Response Due", "date": str(date.today() + timedelta(days=3))}]),
     "links": json.dumps([{"label": "Supplier 8D", "url": "https://example.com/supplier/8d"}]),
     "attachments": json.dumps([])},
    {"id": 4, "topic_group": "Paint & Coating", "sub_topic": "Surface adhesion failure",
     "category": "Elements", "severity": "High",
     "opening_date": datetime.today() - timedelta(days=400), "close_date": None,
     "pic": "K. Vogel", "problem_description": "Peeling after 48h salt-spray test.",
     "root_cause_analysis": "Pre-treatment bath concentration out of range.",
     "corrective_actions": "Bath replenished; batch quarantined.",
     "next_steps": "Retest batch after rework.", "status": "Open", "customer_impact": False,
     "milestone_dates": json.dumps([]), "links": json.dumps([]), "attachments": json.dumps([])},
    {"id": 5, "topic_group": "Welding Defects", "sub_topic": "Undercut on fillet welds",
     "category": "Elements", "severity": "Medium",
     "opening_date": datetime.today() - timedelta(days=45), "close_date": None,
     "pic": "J. Müller", "problem_description": "Undercut exceeding 0.5 mm on fillet joints in zone B.",
     "root_cause_analysis": "Travel speed too high; welder technique issue.",
     "corrective_actions": "Additional welder training completed.",
     "next_steps": "Visual audit of next 5 batches.", "status": "In Progress", "customer_impact": False,
     "milestone_dates": json.dumps([]), "links": json.dumps([]), "attachments": json.dumps([])},
    {"id": 6, "topic_group": "Supplier Quality", "sub_topic": "Late delivery – Component Y7",
     "category": "Cross", "severity": "Medium",
     "opening_date": datetime.today() - timedelta(days=60), "close_date": None,
     "pic": "M. Weber", "problem_description": "Supplier delivering 3–5 days late consistently.",
     "root_cause_analysis": "Raw material shortage at supplier site.",
     "corrective_actions": "Dual-sourcing approval in progress.",
     "next_steps": "Qualify second supplier by end of month.", "status": "Open", "customer_impact": True,
     "milestone_dates": json.dumps([]), "links": json.dumps([]), "attachments": json.dumps([])},
]


def load_data() -> pd.DataFrame:
    if EXCEL_FILE.exists():
        try:
            df = pd.read_excel(EXCEL_FILE, engine="openpyxl")
            df["opening_date"] = pd.to_datetime(df["opening_date"], errors="coerce")
            df["close_date"]   = pd.to_datetime(df.get("close_date"), errors="coerce")
            for col in ["milestone_dates", "links", "attachments"]:
                if col not in df.columns:
                    df[col] = "[]"
                df[col] = df[col].fillna("[]")
            if "severity" not in df.columns:
                df["severity"] = "Medium"
            return df
        except Exception as e:
            st.warning(f"Could not load Excel: {e}. Using demo data.")
    return pd.DataFrame(DEMO_DATA)


def save_data(df: pd.DataFrame):
    _backup()
    out = df.copy()
    out["opening_date"] = out["opening_date"].astype(str)
    out["close_date"]   = out["close_date"].astype(str) if "close_date" in out.columns else ""
    if HAS_OPENPYXL:
        out.to_excel(EXCEL_FILE, index=False, engine="openpyxl")
    else:
        out.to_csv(EXCEL_FILE.with_suffix(".csv"), index=False)


def _backup():
    if EXCEL_FILE.exists():
        ts = datetime.today().strftime("%Y%m%d_%H%M%S")
        shutil.copy(EXCEL_FILE, BACKUP_DIR / f"topics_{ts}.xlsx")


def next_id(df):
    return int(df["id"].max() + 1) if len(df) > 0 else 1


# ── Session state ────────────────────────────────────────────────────────────────
if "df" not in st.session_state:
    st.session_state.df = load_data()
if "edit_id" not in st.session_state:
    st.session_state.edit_id = None

df = enrich(st.session_state.df)

# ── Sidebar ──────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown(
        '<p style="font-family:Syne,sans-serif;font-size:20px;font-weight:800;'
        'color:#F8FAFC;letter-spacing:-.02em;margin-bottom:4px;">🗂️ Topic DB</p>',
        unsafe_allow_html=True,
    )
    st.markdown("---")
    nav = st.radio(
        "Navigation",
        ["📊 Dashboard", "📈 Analytics", "📋 All Topics", "➕ New Topic", "📤 Export"],
        label_visibility="collapsed",
    )
    st.markdown("---")
    st.markdown('<p style="font-size:11px;font-weight:700;text-transform:uppercase;'
                'letter-spacing:.1em;color:#94A3B8;">Filter Topics</p>', unsafe_allow_html=True)
    f_cat    = st.multiselect("Category",  CATEGORIES, default=CATEGORIES)
    f_status = st.multiselect("Status",    STATUSES,   default=STATUSES)
    f_sev    = st.multiselect("Severity",  SEVERITIES, default=SEVERITIES)
    f_search = st.text_input("🔍 Search keywords")
    st.markdown("---")
    st.caption(f"Total topics: **{len(df)}** | Escalated: **{int(df['_escalate'].sum())}**")

# Apply filters
mask = (df["category"].isin(f_cat) & df["status"].isin(f_status) & df["severity"].isin(f_sev))
if f_search:
    q = f_search.lower()
    mask &= (df["topic_group"].str.lower().str.contains(q, na=False) |
             df["sub_topic"].str.lower().str.contains(q, na=False) |
             df["problem_description"].str.lower().str.contains(q, na=False))
dff = df[mask].copy()


# ══════════════════════════════════════════════════════════════════════════════
#  DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
if "Dashboard" in nav:
    st.markdown('<p class="section-title">📊 Summary Dashboard</p>', unsafe_allow_html=True)

    total     = len(df)
    open_n    = len(df[df["status"] == "Open"])
    prog_n    = len(df[df["status"] == "In Progress"])
    blk_n     = len(df[df["status"] == "Blocked"])
    clos_n    = len(df[df["status"] == "Closed"])
    esc_n     = int(df["_escalate"].sum())
    pct       = round(clos_n / total * 100) if total else 0
    avg_risk  = round(df[df["status"] != "Closed"]["_risk"].mean()) if len(df[df["status"] != "Closed"]) else 0

    kpis = [
        ("Total",         total,  "#3B82F6"),
        ("Open",          open_n, "#EF4444"),
        ("In Progress",   prog_n, "#3B82F6"),
        ("Blocked",       blk_n,  "#F59E0B"),
        ("Closed",        clos_n, "#10B981"),
        ("Escalated 🚨",  esc_n,  "#DC2626"),
        (f"Completion",   f"{pct}%", "#6366F1"),
        ("Avg Risk Score",avg_risk, "#8B5CF6"),
    ]
    cols = st.columns(len(kpis))
    for col, (label, val, color) in zip(cols, kpis):
        col.markdown(f"""
        <div class="kpi-card" style="--accent:{color}">
          <div class="kpi-label">{label}</div>
          <div class="kpi-value" style="color:{color}">{val}</div>
        </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # Overall progress bar
    st.markdown(f"""
    <div style="font-size:13px;color:#374151;font-weight:700;font-family:Syne,sans-serif;">
      Overall Completion — {pct}%
    </div>
    <div class="progress-outer" style="height:12px;">
      <div class="progress-inner" style="width:{pct}%;background:linear-gradient(90deg,#3B82F6,#6366F1)"></div>
    </div>
    <div style="font-size:11px;color:#9CA3AF;margin-top:4px;">{clos_n} of {total} topics closed</div>
    """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    col_left, col_right = st.columns(2)

    # ── Aging buckets ──────────────────────────────────────────────────────────
    with col_left:
        st.markdown("**⏱ Aging Buckets (non-closed topics)**")
        active = df[df["status"] != "Closed"].copy()
        b_counts = active["_bucket"].value_counts()
        max_b = max(b_counts.values) if len(b_counts) else 1
        b_colors = {"0–3 Months": "#6EE7B7", "3–6 Months": "#FCD34D", "6–12 Months": "#FCA5A5", "> 1 Year": "#F87171"}
        for b in BUCKET_ORDER:
            cnt   = int(b_counts.get(b, 0))
            pct_b = int(cnt / max_b * 100) if max_b else 0
            color = b_colors.get(b, "#D1D5DB")
            st.markdown(f"""
            <div class="bucket-row">
              <div class="bucket-label">{b}</div>
              <div class="bucket-bar-outer">
                <div class="bucket-bar-inner" style="width:{pct_b}%;background:{color}"></div>
              </div>
              <div class="bucket-count">{cnt}</div>
            </div>""", unsafe_allow_html=True)

    # ── Category breakdown ─────────────────────────────────────────────────────
    with col_right:
        st.markdown("**🏷 Category Breakdown (Pareto)**")
        cat_data = []
        for cat in CATEGORIES:
            sub   = df[df["category"] == cat]
            t     = len(sub)
            cl    = len(sub[sub["status"] == "Closed"])
            cat_data.append({"cat": cat, "total": t, "closed": cl,
                             "pct": round(cl / t * 100) if t else 0})
        cat_data.sort(key=lambda x: x["total"], reverse=True)
        max_t = max(c["total"] for c in cat_data) if cat_data else 1
        for cd in cat_data:
            color  = CAT_HEX.get(cd["cat"], "#6B7280")
            bar_w  = int(cd["total"] / max_t * 100)
            st.markdown(f"""
            <div style="border:1px solid #E5E7EB;border-left:4px solid {color};
                 border-radius:10px;padding:12px 16px;margin-bottom:10px;background:#FAFAFA;">
              <div style="display:flex;justify-content:space-between;align-items:center;">
                <span style="font-weight:700;color:#111827;font-size:14px;font-family:Syne,sans-serif;">{cd['cat']}</span>
                <span style="font-size:12px;color:#6B7280;">{cd['closed']}/{cd['total']} closed ({cd['pct']}%)</span>
              </div>
              <div class="progress-outer" style="margin-top:8px;">
                <div class="progress-inner" style="width:{cd['pct']}%;background:{color}"></div>
              </div>
            </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── High-risk / Escalated agenda list ─────────────────────────────────────
    st.markdown("**🚨 High-Risk & Overdue Agenda List**")
    esc_df = df[df["_escalate"] == True].sort_values("_risk", ascending=False)
    if esc_df.empty:
        st.success("✅ No escalated topics at this time.")
    else:
        for _, row in esc_df.iterrows():
            risk_color = "#EF4444" if row["_risk"] >= 70 else "#F59E0B"
            st.markdown(f"""
            <div class="topic-card" style="--card-accent:{risk_color};border-color:{risk_color}22;">
              <div style="display:flex;justify-content:space-between;align-items:flex-start;gap:12px;">
                <div>
                  <span style="font-family:Syne,sans-serif;font-weight:700;font-size:15px;color:#111827;">
                    #{int(row['id'])} {row['topic_group']}</span>
                  <span style="color:#6B7280;font-size:13px;"> — {row['sub_topic']}</span><br>
                  <span style="font-size:12px;color:#374151;">{row.get('problem_description','')[:100]}…</span>
                </div>
                <div style="text-align:right;flex-shrink:0;">
                  <div class="risk-badge" style="background:{risk_color}22;color:{risk_color};">Risk {row['_risk']}</div><br>
                  <span style="font-size:11px;color:#9CA3AF;">{row['_days_open']}d open | PIC: {row.get('pic','—')}</span>
                </div>
              </div>
            </div>""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
#  ANALYTICS
# ══════════════════════════════════════════════════════════════════════════════
elif "Analytics" in nav:
    st.markdown('<p class="section-title">📈 Analytics</p>', unsafe_allow_html=True)

    # ── Monthly opening trend ──────────────────────────────────────────────────
    st.markdown("**📅 Monthly Opening Trend**")
    trend_df = df.copy()
    trend_df["month"] = pd.to_datetime(trend_df["opening_date"], errors="coerce").dt.to_period("M")
    monthly = trend_df.groupby("month").size().reset_index(name="count")
    monthly["month_str"] = monthly["month"].astype(str)
    monthly = monthly.tail(12)

    if not monthly.empty:
        max_cnt = monthly["count"].max()
        bars_html = ""
        for _, r in monthly.iterrows():
            h = int(r["count"] / max_cnt * 120) if max_cnt else 0
            bars_html += f"""
            <div style="display:flex;flex-direction:column;align-items:center;gap:4px;flex:1;">
              <div style="font-size:11px;font-weight:700;color:#374151;
                          font-family:JetBrains Mono,monospace;">{r['count']}</div>
              <div style="width:100%;max-width:40px;height:{h}px;background:linear-gradient(180deg,#3B82F6,#6366F1);
                          border-radius:4px 4px 0 0;"></div>
              <div style="font-size:10px;color:#9CA3AF;transform:rotate(-30deg);white-space:nowrap;">{r['month_str']}</div>
            </div>"""
        st.markdown(f"""
        <div style="border:1px solid #E5E7EB;border-radius:12px;padding:20px;background:#FAFAFA;">
          <div style="display:flex;align-items:flex-end;gap:8px;height:160px;padding-bottom:24px;">
            {bars_html}
          </div>
        </div>""", unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    col_a, col_b = st.columns(2)

    # ── Root cause Pareto ──────────────────────────────────────────────────────
    with col_a:
        st.markdown("**🔍 Root Cause Pareto (Top Issues by Topic Group)**")
        tg_counts = df[df["status"] != "Closed"]["topic_group"].value_counts().head(10)
        if not tg_counts.empty:
            max_tc = tg_counts.iloc[0]
            for tg, cnt in tg_counts.items():
                bar_w = int(cnt / max_tc * 100)
                st.markdown(f"""
                <div class="bucket-row">
                  <div class="bucket-label" style="width:160px;font-size:11px;">{tg[:22]}</div>
                  <div class="bucket-bar-outer">
                    <div class="bucket-bar-inner" style="width:{bar_w}%;background:#6366F1"></div>
                  </div>
                  <div class="bucket-count">{cnt}</div>
                </div>""", unsafe_allow_html=True)

    # ── Cross-category comparison ──────────────────────────────────────────────
    with col_b:
        st.markdown("**🔀 Cross-Category Comparison**")
        comp_rows = []
        for cat in CATEGORIES:
            sub = df[df["category"] == cat]
            if len(sub) == 0:
                continue
            avg_r = round(sub["_risk"].mean())
            esc_c = int(sub["_escalate"].sum())
            ci_c  = int(sub["customer_impact"].sum())
            avg_d = round(sub["_days_open"].mean())
            comp_rows.append({
                "Category": cat, "Topics": len(sub),
                "Avg Risk": avg_r, "Escalated": esc_c,
                "Customer Impact": ci_c, "Avg Days Open": avg_d,
            })
        if comp_rows:
            cmp_df = pd.DataFrame(comp_rows).set_index("Category")
            st.dataframe(cmp_df, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── Status distribution ────────────────────────────────────────────────────
    st.markdown("**🎯 Status Distribution by Severity**")
    pivot = pd.crosstab(df["severity"], df["status"])
    st.dataframe(pivot, use_container_width=True)


# ══════════════════════════════════════════════════════════════════════════════
#  ALL TOPICS
# ══════════════════════════════════════════════════════════════════════════════
elif "All Topics" in nav:
    st.markdown('<p class="section-title">📋 All Topics</p>', unsafe_allow_html=True)
    st.caption(f"Showing **{len(dff)}** of **{len(df)}** topics")

    sort_col = st.selectbox("Sort by", ["opening_date", "_risk", "_days_open", "status", "category"],
                            format_func=lambda x: x.replace("_","").title())
    sort_asc = st.toggle("Ascending", value=False)

    if dff.empty:
        st.info("No topics match the current filters.")
    else:
        for _, row in dff.sort_values(sort_col, ascending=sort_asc).iterrows():
            cat_color = CAT_HEX.get(row["category"], "#6B7280")
            sta_color = STAT_HEX.get(row["status"], "#6B7280")
            sev_color = SEV_HEX.get(row.get("severity","Medium"), "#FCD34D")
            esc_icon  = "🚨 " if row["_escalate"] else ""

            with st.expander(
                f"{esc_icon}[#{int(row['id'])}] {row['topic_group']} — {row['sub_topic']}  "
                f"|  {row['category']}  |  {row['status']}  |  Risk: {row['_risk']}",
                expanded=False,
            ):
                c1, c2, c3, c4, c5 = st.columns(5)
                c1.metric("Category",    row["category"])
                c2.metric("Status",      row["status"])
                c3.metric("Severity",    row.get("severity","—"))
                c4.metric("PIC",         row.get("pic","—"))
                od = row["opening_date"]
                c5.metric("Days Open",   row["_days_open"])

                col_desc, col_meta = st.columns([3, 1])
                with col_desc:
                    st.markdown(f"**📝 Problem Description**  \n{row.get('problem_description','—')}")
                    st.markdown(f"**🔍 Root Cause Analysis**  \n{row.get('root_cause_analysis','—')}")
                    st.markdown(f"**🔧 Corrective Actions**  \n{row.get('corrective_actions','—')}")
                    st.markdown(f"**🔜 Next Steps**  \n{row.get('next_steps','—')}")

                with col_meta:
                    st.markdown(f"**Opening Date:** {str(od.date()) if not pd.isna(od) else '—'}")
                    st.markdown(f"**Aging Bucket:** {row['_bucket']}")
                    st.markdown(f"**Risk Score:** {row['_risk']}/100")
                    st.markdown(f"**Customer Impact:** {'✅ Yes' if row.get('customer_impact') else '❌ No'}")
                    st.markdown(f"**Escalated:** {'🚨 Yes' if row['_escalate'] else '✅ No'}")

                    # Milestone dates
                    try:
                        milestones = json.loads(row.get("milestone_dates") or "[]")
                        if milestones:
                            st.markdown("**📅 Milestones:**")
                            for m in milestones:
                                st.markdown(f"- {m['label']}: `{m['date']}`")
                    except Exception:
                        pass

                    # Links
                    try:
                        links = json.loads(row.get("links") or "[]")
                        if links:
                            st.markdown("**🔗 Links:**")
                            for lnk in links:
                                st.markdown(f"- [{lnk['label']}]({lnk['url']})")
                    except Exception:
                        pass

                    # Attachments
                    try:
                        attachments = json.loads(row.get("attachments") or "[]")
                        if attachments:
                            st.markdown("**📎 Attachments:**")
                            for att in attachments:
                                att_path = ATTACH_DIR / att
                                if att_path.exists():
                                    with open(att_path, "rb") as fh:
                                        st.download_button(
                                            label=f"⬇️ {att}", data=fh,
                                            file_name=att, key=f"dl_{row['id']}_{att}"
                                        )
                                else:
                                    st.caption(f"📎 {att} (not found)")
                    except Exception:
                        pass

                if st.button(f"✏️ Edit Topic #{int(row['id'])}", key=f"edit_{row['id']}"):
                    st.session_state.edit_id = int(row["id"])
                    st.rerun()

    # ── Inline edit form ───────────────────────────────────────────────────────
    if st.session_state.edit_id is not None:
        eid = st.session_state.edit_id
        base_df = st.session_state.df
        row_match = base_df[base_df["id"] == eid]
        if len(row_match) > 0:
            row = row_match.iloc[0]
            st.markdown("---")
            st.markdown(f"### ✏️ Edit Topic #{eid}")
            with st.form(f"edit_form_{eid}"):
                col1, col2 = st.columns(2)
                tg   = col1.text_input("Topic Group",    value=row["topic_group"])
                st_  = col1.text_input("Sub-Topic",      value=row["sub_topic"])
                cat  = col1.selectbox("Category",  CATEGORIES, index=CATEGORIES.index(row["category"]) if row["category"] in CATEGORIES else 0)
                sev  = col1.selectbox("Severity",  SEVERITIES, index=SEVERITIES.index(row.get("severity","Medium")) if row.get("severity","Medium") in SEVERITIES else 1)
                pic  = col2.text_input("PIC",            value=row.get("pic",""))
                stat = col2.selectbox("Status",    STATUSES,   index=STATUSES.index(row["status"]) if row["status"] in STATUSES else 0)
                od   = col2.date_input("Opening Date", value=row["opening_date"].date() if not pd.isna(row["opening_date"]) else date.today())
                cd_val = row.get("close_date")
                if stat == "Closed":
                    cd = col2.date_input("Close Date", value=cd_val.date() if cd_val and not pd.isna(cd_val) else date.today())
                else:
                    cd = None
                ci   = col2.checkbox("Customer Impact", value=bool(row.get("customer_impact", False)))

                pd_  = st.text_area("Problem Description",  value=row.get("problem_description",""), height=80)
                rca  = st.text_area("Root Cause Analysis",  value=row.get("root_cause_analysis",""),  height=80)
                ca   = st.text_area("Corrective Actions",   value=row.get("corrective_actions",""),   height=80)
                ns   = st.text_area("Next Steps",            value=row.get("next_steps",""),           height=60)

                links_raw = st.text_area(
                    "Links (JSON) — format: [{\"label\":\"8D\",\"url\":\"https://…\"}]",
                    value=row.get("links","[]"), height=60
                )
                milestones_raw = st.text_area(
                    "Milestone Dates (JSON) — format: [{\"label\":\"8D Due\",\"date\":\"2025-12-31\"}]",
                    value=row.get("milestone_dates","[]"), height=60
                )

                col_a, col_b = st.columns(2)
                save   = col_a.form_submit_button("💾 Save Changes", type="primary")
                cancel = col_b.form_submit_button("❌ Cancel")

                if save:
                    idx = base_df.index[base_df["id"] == eid][0]
                    for field, val in [
                        ("topic_group", tg), ("sub_topic", st_), ("category", cat),
                        ("severity", sev), ("pic", pic), ("status", stat),
                        ("opening_date", pd.Timestamp(od)),
                        ("close_date", pd.Timestamp(cd) if cd else pd.NaT),
                        ("customer_impact", ci),
                        ("problem_description", pd_), ("root_cause_analysis", rca),
                        ("corrective_actions", ca), ("next_steps", ns),
                        ("links", links_raw), ("milestone_dates", milestones_raw),
                    ]:
                        base_df.at[idx, field] = val
                    st.session_state.df = base_df
                    save_data(base_df)
                    st.session_state.edit_id = None
                    st.success("✅ Topic updated!")
                    st.rerun()
                if cancel:
                    st.session_state.edit_id = None
                    st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
#  NEW TOPIC
# ══════════════════════════════════════════════════════════════════════════════
elif "New Topic" in nav:
    st.markdown('<p class="section-title">➕ New Topic</p>', unsafe_allow_html=True)

    # Auto-fill defaults
    st.info("Fields marked * are required. TopicID, DaysOpen, Aging Bucket, Risk Score and Escalation Flag are generated automatically.")

    with st.form("new_topic_form"):
        col1, col2 = st.columns(2)
        tg   = col1.text_input("Topic Group *")
        st_  = col1.text_input("Sub-Topic *")
        cat  = col1.selectbox("Category *", CATEGORIES)
        sev  = col1.selectbox("Severity *", SEVERITIES, index=1)

        pic  = col2.text_input("Person in Charge (PIC)")
        stat = col2.selectbox("Status *", STATUSES)
        od   = col2.date_input("Opening Date", value=date.today())
        if stat == "Closed":
            cd = col2.date_input("Close Date", value=date.today())
        else:
            cd = None
        ci   = col2.checkbox("Customer Impact")

        pd_  = st.text_area("Problem Description *", height=80)
        rca  = st.text_area("Root Cause Analysis",   height=80)
        ca   = st.text_area("Corrective Actions",    height=80)
        ns   = st.text_area("Next Steps",             height=60)

        st.markdown("**🔗 Links** (one per line, format: `Label | URL`)")
        links_input = st.text_area("Links", height=80, placeholder="8D Report | https://…\nMeeting Notes | https://…")

        st.markdown("**📅 Milestone Dates** (one per line, format: `Label | YYYY-MM-DD`)")
        milestones_input = st.text_area("Milestones", height=60, placeholder="8D Due | 2025-12-31\nAudit Date | 2026-01-15")

        uploaded_files = st.file_uploader(
            "📎 Attach Files (photos, documents — PDF, PNG, JPG, DOCX, XLSX)",
            accept_multiple_files=True,
            type=["png", "jpg", "jpeg", "pdf", "docx", "xlsx", "txt"]
        )

        submitted = st.form_submit_button("✅ Save Topic", type="primary")
        if submitted:
            if not tg or not st_ or not pd_:
                st.error("Please fill in Topic Group, Sub-Topic and Problem Description.")
            else:
                new_id = next_id(st.session_state.df)

                # Parse links
                links_list = []
                for line in links_input.strip().splitlines():
                    if "|" in line:
                        parts = line.split("|", 1)
                        links_list.append({"label": parts[0].strip(), "url": parts[1].strip()})

                # Parse milestones
                milestones_list = []
                for line in milestones_input.strip().splitlines():
                    if "|" in line:
                        parts = line.split("|", 1)
                        milestones_list.append({"label": parts[0].strip(), "date": parts[1].strip()})

                # Save attachments
                saved_files = []
                if uploaded_files:
                    for f in uploaded_files:
                        safe_name = sanitize_filename(f"{new_id}_{f.name}")
                        dest      = ATTACH_DIR / safe_name
                        dest.write_bytes(f.getvalue())
                        saved_files.append(safe_name)

                new_row = {
                    "id": new_id, "topic_group": tg, "sub_topic": st_,
                    "category": cat, "severity": sev,
                    "opening_date": pd.Timestamp(od),
                    "close_date": pd.Timestamp(cd) if cd else pd.NaT,
                    "pic": pic, "problem_description": pd_,
                    "root_cause_analysis": rca, "corrective_actions": ca,
                    "next_steps": ns, "status": stat, "customer_impact": ci,
                    "links": json.dumps(links_list),
                    "milestone_dates": json.dumps(milestones_list),
                    "attachments": json.dumps(saved_files),
                }
                new_df = pd.concat([st.session_state.df, pd.DataFrame([new_row])], ignore_index=True)
                st.session_state.df = new_df
                save_data(new_df)
                st.success(f"✅ Topic #{new_id} saved! (Risk will be calculated automatically)")
                st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
#  EXPORT
# ══════════════════════════════════════════════════════════════════════════════
elif "Export" in nav:
    st.markdown('<p class="section-title">📤 Export</p>', unsafe_allow_html=True)

    export_all = st.toggle("Export ALL topics (ignore filters)", value=True)
    export_df  = enrich(st.session_state.df) if export_all else dff.copy()
    st.info(f"Will export **{len(export_df)}** topics.")

    col_ex1, col_ex2, col_ex3 = st.columns(3)

    # ── Excel ──────────────────────────────────────────────────────────────────
    with col_ex1:
        st.subheader("📊 Excel")

        def build_excel(data: pd.DataFrame) -> bytes:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            HDR  = Font(name="Calibri", size=13, bold=True, color="FFFFFF")
            HDR2 = Font(name="Calibri", size=11, bold=True, color="1F3864")
            NORM = Font(name="Calibri", size=10)
            BOLD = Font(name="Calibri", size=10, bold=True)
            BLUE_F = PatternFill("solid", fgColor="1F3864")
            MED_F  = PatternFill("solid", fgColor="2E75B6")
            LBL_F  = PatternFill("solid", fgColor="D9E1F2")
            CAT_F  = {"Elements": PatternFill("solid", fgColor="DBEAFE"),
                      "Assembly": PatternFill("solid", fgColor="D1FAE5"),
                      "Cross":    PatternFill("solid", fgColor="FEF3C7")}
            STA_F  = {"Open":        PatternFill("solid", fgColor="FEE2E2"),
                      "In Progress": PatternFill("solid", fgColor="DBEAFE"),
                      "Blocked":     PatternFill("solid", fgColor="FEF3C7"),
                      "Closed":      PatternFill("solid", fgColor="D1FAE5")}
            thin = Side(style="thin", color="B0B0B0")
            brd  = Border(left=thin, right=thin, top=thin, bottom=thin)

            def hcell(ws, r, text, cols=10):
                ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=cols)
                c = ws.cell(row=r, column=1, value=text)
                c.font = HDR; c.fill = BLUE_F
                c.alignment = Alignment(horizontal="center", vertical="center")
                ws.row_dimensions[r].height = 22

            def wcell(ws, r, col, val, bold=False, fill=None, align="left"):
                c = ws.cell(row=r, column=col, value=val)
                c.font = BOLD if bold else NORM
                if fill: c.fill = fill
                c.alignment = Alignment(horizontal=align, vertical="center", wrap_text=True)
                c.border = brd
                return c

            # Sheet 1: Summary
            ws = wb.active; ws.title = "Summary"
            hcell(ws, 1, "TOPIC DATABASE — SUMMARY DASHBOARD", cols=8)
            total_n  = len(data)
            closed_n = len(data[data["status"] == "Closed"])
            kpi_labels = ["Total","Open","In Progress","Blocked","Closed","Escalated","Completion %","Avg Risk"]
            kpi_vals   = [
                total_n,
                len(data[data["status"]=="Open"]),
                len(data[data["status"]=="In Progress"]),
                len(data[data["status"]=="Blocked"]),
                closed_n,
                int(data["_escalate"].sum()),
                f"{round(closed_n/total_n*100) if total_n else 0}%",
                round(data[data["status"]!="Closed"]["_risk"].mean()) if len(data[data["status"]!="Closed"])>0 else 0,
            ]
            ws.cell(row=3, column=1, value="KPI Overview").font = HDR2
            for i, (l, v) in enumerate(zip(kpi_labels, kpi_vals)):
                wcell(ws, 4, i+1, l, bold=True, fill=LBL_F, align="center")
                wcell(ws, 5, i+1, v, align="center")

            # Category table
            ws.cell(row=7, column=1, value="Category Breakdown").font = HDR2
            for i, h in enumerate(["Category","Total","Open","In Progress","Blocked","Closed","% Done","Escalated"]):
                wcell(ws, 8, i+1, h, bold=True, fill=LBL_F, align="center")
            for ri, cat in enumerate(CATEGORIES):
                sub = data[data["category"]==cat]; t=len(sub); cl=len(sub[sub["status"]=="Closed"])
                row_n = 9+ri
                vals  = [cat, t, len(sub[sub["status"]=="Open"]),
                         len(sub[sub["status"]=="In Progress"]),
                         len(sub[sub["status"]=="Blocked"]), cl,
                         f"{round(cl/t*100) if t else 0}%",
                         int(sub["_escalate"].sum())]
                for ci_, v in enumerate(vals):
                    wcell(ws, row_n, ci_+1, v, align="center").fill = CAT_F.get(cat, PatternFill())

            for i, w in enumerate([22,12,14,14,12,12,12,12], 1):
                ws.column_dimensions[get_column_letter(i)].width = w

            # Sheet 2: Topic Documentation
            ws2 = wb.create_sheet("Topic Documentation")
            cols2 = [
                ("ID",6),("Topic Group",22),("Sub-Topic",25),("Category",14),
                ("Severity",12),("Opening Date",14),("PIC",16),("Status",14),
                ("Days Open",12),("Risk Score",12),("Escalated",12),
                ("Problem Description",42),("Root Cause",40),("Corrective Actions",40),
                ("Next Steps",30),("Customer Impact",16),("Aging Bucket",16),
            ]
            hcell(ws2, 1, "TOPIC DOCUMENTATION", cols=len(cols2))
            for ci_, (name, width) in enumerate(cols2, 1):
                c = ws2.cell(row=2, column=ci_, value=name)
                c.font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
                c.fill = MED_F
                c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                c.border = brd
                ws2.column_dimensions[get_column_letter(ci_)].width = width

            for ri, (_, r_) in enumerate(data.sort_values("id").iterrows(), start=3):
                od = r_["opening_date"]
                vals2 = [
                    int(r_["id"]), r_["topic_group"], r_["sub_topic"], r_["category"],
                    r_.get("severity","—"),
                    od.date() if not pd.isna(od) else "",
                    r_.get("pic",""), r_["status"],
                    int(r_["_days_open"]), int(r_["_risk"]),
                    "YES" if r_["_escalate"] else "no",
                    r_.get("problem_description",""), r_.get("root_cause_analysis",""),
                    r_.get("corrective_actions",""), r_.get("next_steps",""),
                    "Yes" if r_.get("customer_impact") else "No", r_["_bucket"],
                ]
                for ci_, val in enumerate(vals2, 1):
                    c = ws2.cell(row=ri, column=ci_, value=val)
                    c.font = NORM
                    c.alignment = Alignment(vertical="top", wrap_text=True)
                    c.border = brd
                ws2.cell(row=ri, column=4).fill  = CAT_F.get(r_["category"], PatternFill())
                ws2.cell(row=ri, column=8).fill  = STA_F.get(r_["status"],   PatternFill())
                ws2.row_dimensions[ri].height    = 42
            ws2.freeze_panes = "A3"

            # Sheet 3: Escalated
            ws3 = wb.create_sheet("Escalated Topics")
            esc = data[data["_escalate"] == True].sort_values("_risk", ascending=False)
            hcell(ws3, 1, "ESCALATED TOPICS", cols=8)
            esc_hdr = ["ID","Topic Group","Sub-Topic","Status","Risk Score","Days Open","PIC","Problem"]
            for ci_, h in enumerate(esc_hdr, 1):
                c = ws3.cell(row=2, column=ci_, value=h)
                c.font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
                c.fill = PatternFill("solid", fgColor="DC2626")
                c.alignment = Alignment(horizontal="center", vertical="center")
                c.border = brd
                ws3.column_dimensions[get_column_letter(ci_)].width = [6,22,25,14,12,12,16,40][ci_-1]
            for ri, (_, r_) in enumerate(esc.iterrows(), start=3):
                for ci_, val in enumerate([
                    int(r_["id"]), r_["topic_group"], r_["sub_topic"], r_["status"],
                    int(r_["_risk"]), int(r_["_days_open"]), r_.get("pic",""),
                    r_.get("problem_description","")[:120],
                ], 1):
                    c = ws3.cell(row=ri, column=ci_, value=val)
                    c.font = NORM; c.border = brd
                    c.alignment = Alignment(vertical="top", wrap_text=True)
            ws3.freeze_panes = "A3"

            buf = io.BytesIO(); wb.save(buf); return buf.getvalue()

        if HAS_OPENPYXL:
            xl = build_excel(export_df)
            st.download_button("⬇️ Download Excel", data=xl,
                               file_name=f"topic_report_{datetime.today().strftime('%Y%m%d')}.xlsx",
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               type="primary")
        else:
            st.warning("Install openpyxl: `pip install openpyxl`")

    # ── CSV ────────────────────────────────────────────────────────────────────
    with col_ex2:
        st.subheader("📄 CSV")
        csv_df = export_df.copy()
        csv_df["opening_date"] = csv_df["opening_date"].astype(str)
        csv_df["close_date"]   = csv_df.get("close_date", pd.NaT).astype(str)
        csv_bytes = csv_df.to_csv(index=False).encode("utf-8-sig")
        st.download_button("⬇️ Download CSV", data=csv_bytes,
                           file_name=f"topic_report_{datetime.today().strftime('%Y%m%d')}.csv",
                           mime="text/csv", type="primary")

    # ── PPTX ──────────────────────────────────────────────────────────────────
    with col_ex3:
        st.subheader("📑 PowerPoint")

        def build_pptx(data: pd.DataFrame) -> bytes:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]  # blank

            def rgb(h): r,g,b = int(h[1:3],16),int(h[3:5],16),int(h[5:7],16); return RGBColor(r,g,b)
            DARK = rgb("#0F172A"); BLUE = rgb("#3B82F6"); GRAY = rgb("#6B7280"); WHITE = rgb("#FFFFFF")
            CAT_RGB  = {"Elements": rgb("#3B82F6"), "Assembly": rgb("#10B981"), "Cross": rgb("#F59E0B")}
            STAT_RGB = {"Open": rgb("#EF4444"), "In Progress": rgb("#3B82F6"), "Blocked": rgb("#F59E0B"), "Closed": rgb("#10B981")}

            def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
                txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf  = txb.text_frame; tf.word_wrap = True
                p   = tf.paragraphs[0]; p.text = text; p.alignment = align
                run = p.runs[0]; run.font.size = Pt(size); run.font.bold = bold
                if color: run.font.color.rgb = color
                return txb

            def add_rect(slide, left, top, width, height, fill_color, line_color=None):
                from pptx.util import Inches as I_
                shape = slide.shapes.add_shape(1, I_(left), I_(top), I_(width), I_(height))
                shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
                if line_color: shape.line.color.rgb = line_color
                else: shape.line.fill.background()
                return shape

            # ── Title slide ──────────────────────────────────────────────────
            slide = prs.slides.add_slide(blank_layout)
            add_rect(slide, 0, 0, 13.33, 7.5, DARK)
            add_rect(slide, 0, 5.5, 13.33, 2.0, BLUE)
            add_textbox(slide, 0.5, 1.5, 12, 1.5, "Topic Database", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_textbox(slide, 0.5, 3.2, 12, 0.6,
                        f"Report generated: {datetime.today().strftime('%Y-%m-%d %H:%M')}  |  Total topics: {len(data)}",
                        size=14, color=rgb("#94A3B8"), align=PP_ALIGN.CENTER)
            total_n  = len(data); closed_n = len(data[data["status"]=="Closed"])
            pct_done = round(closed_n/total_n*100) if total_n else 0
            add_textbox(slide, 0.5, 5.7, 12, 0.8,
                        f"Open: {len(data[data['status']=='Open'])}   In Progress: {len(data[data['status']=='In Progress'])}   "
                        f"Blocked: {len(data[data['status']=='Blocked'])}   Closed: {closed_n}   Completion: {pct_done}%",
                        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

            # ── Summary KPI slide ─────────────────────────────────────────────
            slide = prs.slides.add_slide(blank_layout)
            add_rect(slide, 0, 0, 13.33, 1.0, DARK)
            add_textbox(slide, 0.3, 0.1, 12, 0.8, "Summary Dashboard", size=22, bold=True, color=WHITE)
            kpis = [("Total", len(data)), ("Open", len(data[data["status"]=="Open"])),
                    ("In Progress", len(data[data["status"]=="In Progress"])),
                    ("Blocked", len(data[data["status"]=="Blocked"])),
                    ("Closed", closed_n), ("Completion", f"{pct_done}%")]
            for i, (lbl, val) in enumerate(kpis):
                x = 0.3 + i * 2.15
                add_rect(slide, x, 1.3, 2.0, 1.4, rgb("#F8FAFC"))
                add_textbox(slide, x+0.1, 1.35, 1.8, 0.5, lbl, size=9, color=GRAY, align=PP_ALIGN.CENTER)
                add_textbox(slide, x+0.1, 1.75, 1.8, 0.7, str(val), size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)

            # ── Per-topic slides ─────────────────────────────────────────────
            for _, row in data.sort_values("id").iterrows():
                if row["status"] == "Closed":
                    continue  # Optionally skip closed – remove this line to include all
                slide = prs.slides.add_slide(blank_layout)
                cat   = row["category"]
                stat  = row["status"]
                c_col = CAT_RGB.get(cat, BLUE)
                s_col = STAT_RGB.get(stat, GRAY)

                # Header bar
                add_rect(slide, 0, 0, 13.33, 1.1, DARK)
                add_textbox(slide, 0.3, 0.08, 10, 0.55,
                            f"#{int(row['id'])}  {row['topic_group']} — {row['sub_topic']}",
                            size=18, bold=True, color=WHITE)
                add_textbox(slide, 0.3, 0.6, 4, 0.4,
                            f"Category: {cat}  |  Status: {stat}  |  PIC: {row.get('pic','—')}",
                            size=10, color=rgb("#94A3B8"))
                add_textbox(slide, 8.0, 0.08, 5, 0.5,
                            f"Risk: {int(row['_risk'])}/100  |  {int(row['_days_open'])}d open",
                            size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

                # Accent bar
                add_rect(slide, 0, 1.1, 13.33, 0.07, c_col)

                # Content columns
                col_data = [
                    ("📝 Problem Description", row.get("problem_description","—")),
                    ("🔍 Root Cause Analysis", row.get("root_cause_analysis","—")),
                    ("🔧 Corrective Actions",  row.get("corrective_actions","—")),
                    ("🔜 Next Steps",           row.get("next_steps","—") or "—"),
                ]
                for i, (title, content) in enumerate(col_data):
                    x = 0.2 + i * 3.3
                    add_rect(slide, x, 1.3, 3.1, 4.5, rgb("#F8FAFC"))
                    add_textbox(slide, x+0.1, 1.35, 2.9, 0.45, title, size=9, bold=True, color=c_col)
                    add_textbox(slide, x+0.1, 1.8, 2.9, 3.8, (content or "—")[:300], size=8.5, color=DARK)

                # Footer
                od = row["opening_date"]
                add_textbox(slide, 0.2, 6.2, 13, 0.4,
                            f"Opened: {str(od.date()) if not pd.isna(od) else '—'}  |  "
                            f"Aging: {row['_bucket']}  |  "
                            f"Customer Impact: {'YES' if row.get('customer_impact') else 'No'}  |  "
                            f"Severity: {row.get('severity','—')}  |  "
                            f"{'🚨 ESCALATED' if row['_escalate'] else ''}",
                            size=9, color=GRAY)

            buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

        if HAS_PPTX:
            pptx_bytes = build_pptx(export_df)
            st.download_button("⬇️ Download PPTX", data=pptx_bytes,
                               file_name=f"topic_report_{datetime.today().strftime('%Y%m%d')}.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                               type="primary")
        else:
            st.warning("Install python-pptx: `pip install python-pptx`")
            st.code("pip install python-pptx")

    st.markdown("---")
    st.markdown("**🗄 Data Storage Info**")
    st.json({
        "Excel file": str(EXCEL_FILE),
        "Attachments": str(ATTACH_DIR),
        "Backups": str(BACKUP_DIR),
        "Total topics": len(st.session_state.df),
        "Backup count": len(list(BACKUP_DIR.glob("*.xlsx"))),
    })
